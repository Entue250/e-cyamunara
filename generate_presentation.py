from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Brand colors ──────────────────────────────────────────────────────────────
RNP_BLUE  = RGBColor(0x00, 0x30, 0x87)
GOLD      = RGBColor(0xFF, 0xD7, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFF)
GOLD_LIGHT= RGBColor(0xFF, 0xF0, 0x80)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

SLIDE_NUM = [0]   # mutable counter

# ── Helper: solid fill ─────────────────────────────────────────────────────────
def solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

# ── Helper: add text box ───────────────────────────────────────────────────────
def add_text(slide, text, x, y, w, h, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

# ── Helper: add rectangle ─────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    solid_fill(shape, color)
    shape.line.fill.background()
    return shape

# ── Helper: blue header bar ───────────────────────────────────────────────────
def header_bar(slide, title_text, subtitle=None):
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), RNP_BLUE)
    add_text(slide, title_text,
             Inches(0.4), Inches(0.12), Inches(11.5), Inches(0.7),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.78), Inches(11.5), Inches(0.42),
                 size=14, color=GOLD)

# ── Helper: gold accent line ──────────────────────────────────────────────────
def gold_line(slide, y=Inches(1.35)):
    add_rect(slide, 0, y, SLIDE_W, Inches(0.045), GOLD)

# ── Helper: footer ────────────────────────────────────────────────────────────
def footer(slide, num):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), RNP_BLUE)
    add_text(slide, "RNP E-CYAMUNARA  |  Online Vehicle Auction Platform",
             Inches(0.3), Inches(7.16), Inches(10), Inches(0.3),
             size=9, color=GOLD)
    add_text(slide, str(num),
             Inches(12.6), Inches(7.16), Inches(0.6), Inches(0.3),
             size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# ── Helper: bullet list ───────────────────────────────────────────────────────
def bullet_list(slide, items, x, y, w, h, size=16, dot_color=GOLD):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        # bullet dot
        dot = p.add_run()
        dot.text = "▸  "
        dot.font.size = Pt(size - 2)
        dot.font.color.rgb = dot_color
        dot.font.bold = True
        # text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = DARK_TEXT

# ── Helper: two-column bullets ────────────────────────────────────────────────
def two_col_bullets(slide, left_items, right_items, y=Inches(1.6),
                    col_h=Inches(4.8), size=15):
    bullet_list(slide, left_items,  Inches(0.4), y, Inches(5.8), col_h, size)
    bullet_list(slide, right_items, Inches(6.5), y, Inches(6.2), col_h, size)

# ── Helper: section card (colored box) ───────────────────────────────────────
def section_card(slide, title, items, x, y, w, h, bg_color=LIGHT_BG):
    card = add_rect(slide, x, y, w, h, bg_color)
    add_text(slide, title,
             x + Inches(0.18), y + Inches(0.1), w - Inches(0.3), Inches(0.38),
             size=14, bold=True, color=RNP_BLUE)
    line = add_rect(slide, x + Inches(0.12), y + Inches(0.5),
                    w - Inches(0.24), Inches(0.03), GOLD)
    bullet_list(slide, items,
                x + Inches(0.12), y + Inches(0.58),
                w - Inches(0.24), h - Inches(0.7),
                size=12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
def slide_title():
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    SLIDE_NUM[0] += 1

    # Full blue background top half
    add_rect(sl, 0, 0, SLIDE_W, Inches(4.4), RNP_BLUE)

    # Gold decorative stripe
    add_rect(sl, 0, Inches(4.4), SLIDE_W, Inches(0.12), GOLD)

    # Project name
    add_text(sl, "RNP E-CYAMUNARA",
             Inches(0.6), Inches(0.55), Inches(12), Inches(1.1),
             size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "ONLINE AUCTION PLATFORM",
             Inches(0.6), Inches(1.55), Inches(12), Inches(0.7),
             size=26, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(sl, "Digital Vehicle Auction Management for Rwanda National Police",
             Inches(0.6), Inches(2.25), Inches(12), Inches(0.55),
             size=15, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Module & group — white area
    add_text(sl, "Mobile Application and Handheld Devices  |  GROUP 6",
             Inches(0.6), Inches(3.15), Inches(12), Inches(0.45),
             size=13, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

    # Team members section
    add_rect(sl, Inches(0.5), Inches(4.7), Inches(5.5), Inches(2.2), LIGHT_BG)
    add_text(sl, "Team Members",
             Inches(0.7), Inches(4.78), Inches(5), Inches(0.35),
             size=13, bold=True, color=RNP_BLUE)
    members = [
        "IRUMVA Jean Marie",
        "MANZI Stiven",
        "NDAYISHIMIYE Issa",
        "NIYOMUGABO Eduard",
    ]
    bullet_list(sl, members, Inches(0.7), Inches(5.15), Inches(5), Inches(1.55), size=13)

    # Date / academic year
    add_rect(sl, Inches(7.0), Inches(4.7), Inches(5.7), Inches(2.2), LIGHT_BG)
    add_text(sl, "Presentation Details",
             Inches(7.2), Inches(4.78), Inches(5.2), Inches(0.35),
             size=13, bold=True, color=RNP_BLUE)
    add_text(sl, "Academic Year:  2024 – 2025",
             Inches(7.2), Inches(5.18), Inches(5.2), Inches(0.38), size=13, color=DARK_TEXT)
    add_text(sl, "Date:  May 2025",
             Inches(7.2), Inches(5.6),  Inches(5.2), Inches(0.38), size=13, color=DARK_TEXT)
    add_text(sl, "Institution:  RP – Musanze Campus",
             Inches(7.2), Inches(6.0),  Inches(5.2), Inches(0.38), size=13, color=DARK_TEXT)

    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION
# ══════════════════════════════════════════════════════════════════════════════
def slide_introduction():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Introduction", "What is E-CYAMUNARA?")
    gold_line(sl)

    add_text(sl, "E-CYAMUNARA is a mobile-first online auction platform built for the Rwanda National Police (RNP) to digitize and modernize the management of vehicle auctions.",
             Inches(0.4), Inches(1.55), Inches(12.4), Inches(0.85),
             size=15, color=DARK_TEXT)

    items_l = [
        "Purpose-built Flutter mobile application",
        "Replaces manual, paper-based auction processes",
        "Enables citizens to participate in RNP auctions remotely",
        "Real-time bidding with live auction updates",
    ]
    items_r = [
        "Covers full auction lifecycle: post → bid → close → report",
        "Three-tier role system for controlled access",
        "Secure data handling with Supabase backend",
        "Developed as final-year capstone project — Group 6",
    ]
    two_col_bullets(sl, items_l, items_r, y=Inches(2.55))
    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
def slide_problem():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Problem Statement", "Challenges with traditional vehicle auctions")
    gold_line(sl)

    problems = [
        ("Physical Attendance Required",
         ["Citizens must travel to auction venues", "High transport cost and time burden"]),
        ("Limited Transparency",
         ["No public view of bid history", "Fraud risk in manual bid recording"]),
        ("Delayed Communication",
         ["Winners notified only in person or by phone", "No automated status updates"]),
        ("Poor Record Management",
         ["Paper-based documentation lost or damaged", "No audit trail"]),
        ("Limited Accessibility",
         ["Auctions inaccessible to citizens outside the region", "Language and scheduling barriers"]),
        ("Difficult Bid Tracking",
         ["No system to track personal bids or outbid status", "Buyers must call RNP for updates"]),
    ]

    # 3-column card layout
    positions = [
        (Inches(0.3),  Inches(1.55), Inches(3.9), Inches(2.3)),
        (Inches(4.5),  Inches(1.55), Inches(3.9), Inches(2.3)),
        (Inches(8.7),  Inches(1.55), Inches(4.2), Inches(2.3)),
        (Inches(0.3),  Inches(4.05), Inches(3.9), Inches(2.3)),
        (Inches(4.5),  Inches(4.05), Inches(3.9), Inches(2.3)),
        (Inches(8.7),  Inches(4.05), Inches(4.2), Inches(2.3)),
    ]
    for i, (title, pts) in enumerate(problems):
        x, y, w, h = positions[i]
        section_card(sl, title, pts, x, y, w, h)

    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROPOSED SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
def slide_solution():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Proposed Solution", "How E-CYAMUNARA solves the problem")
    gold_line(sl)

    solutions = [
        "Online Auction Management — RNP admins post, edit, and close auctions from any device",
        "Real-Time Bidding — Supabase Realtime streams live bid updates instantly to all participants",
        "Mobile Accessibility — Android app built with Flutter; accessible from anywhere in Rwanda",
        "Push Notifications — OneSignal alerts users when outbid, when auctions open/close",
        "Digital Reports — Auto-generated region and national PDF auction reports",
        "Transparent Auction Process — All bids are logged, timestamped, and immutable",
        "Secure RBAC — Role-based access ensures only authorized users perform sensitive actions",
        "Centralized Administration — Super admin oversees all regions and admins from one dashboard",
    ]
    bullet_list(sl, solutions, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.3), size=15)
    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM USERS / ROLES
# ══════════════════════════════════════════════════════════════════════════════
def slide_roles():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "System Users & Roles", "Three-tier role-based access control")
    gold_line(sl)

    roles = [
        ("CLIENT", [
            "Registers and logs in via mobile app",
            "Selects preferred auction region",
            "Browses active auction listings",
            "Places and tracks bids in real time",
            "Receives outbid and win notifications",
            "Submits feedback on completed auctions",
            "Views personal bid history",
        ]),
        ("REGION ADMIN", [
            "Manages auctions for assigned RNP region",
            "Posts new vehicle auctions with photos",
            "Monitors all bids on active auctions",
            "Manually closes auctions and declares winner",
            "Manages client accounts (suspend / activate)",
            "Generates region-level auction reports",
            "Views client feedback and analytics",
        ]),
        ("SUPER ADMIN", [
            "Full oversight of all regions nationwide",
            "Creates and manages region admin accounts",
            "Monitors all auctions across all regions",
            "Accesses national-level reports and analytics",
            "Views system-wide statistics",
            "Has ultimate authority over platform settings",
        ]),
    ]

    xs = [Inches(0.25), Inches(4.55), Inches(8.85)]
    for i, (role_name, pts) in enumerate(roles):
        x = xs[i]
        w = Inches(4.0)
        h = Inches(4.85)
        y = Inches(1.55)
        card = add_rect(sl, x, y, w, h, LIGHT_BG)
        add_rect(sl, x, y, w, Inches(0.52), RNP_BLUE)
        add_text(sl, role_name,
                 x + Inches(0.1), y + Inches(0.07), w - Inches(0.2), Inches(0.38),
                 size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        bullet_list(sl, pts, x + Inches(0.15), y + Inches(0.6),
                    w - Inches(0.25), h - Inches(0.7), size=13)

    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MAIN FEATURES
# ══════════════════════════════════════════════════════════════════════════════
def slide_features():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Main Features", "Core capabilities of the platform")
    gold_line(sl)

    feature_blocks = [
        ("Authentication & RBAC",     ["JWT-based Supabase Auth", "Role-resolved routing with GoRouter"]),
        ("Auction Management",        ["Post, edit, monitor, close auctions", "Photo upload via Supabase Storage"]),
        ("Real-Time Bidding",         ["Live bid streaming via Supabase Realtime", "Atomic bid validation via Edge Functions"]),
        ("Push Notifications",        ["OneSignal integration", "Outbid and auction-close alerts"]),
        ("Reports & Analytics",       ["Region and national PDF reports", "Dynamic stat counters on dashboards"]),
        ("Client Management",         ["Suspend / activate accounts", "Feedback review per client"]),
        ("Feedback System",           ["Clients rate completed auctions", "Admins view all submitted feedback"]),
        ("Region-Based Filtering",    ["Clients select their region on signup", "Admins scoped to assigned region"]),
        ("Secure Data Handling",      ["National IDs encrypted with AES-256", "RLS policies enforce row-level security"]),
    ]

    xs = [Inches(0.25), Inches(4.55), Inches(8.85)]
    ys = [Inches(1.58), Inches(3.25), Inches(4.92)]
    w  = Inches(4.0)
    h  = Inches(1.5)

    for i, (title, pts) in enumerate(feature_blocks):
        row = i // 3
        col = i  % 3
        x = xs[col]
        y = ys[row]
        section_card(sl, title, pts, x, y, w, h)

    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TECHNOLOGIES USED
# ══════════════════════════════════════════════════════════════════════════════
def slide_tech():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Technologies Used", "Stack powering E-CYAMUNARA")
    gold_line(sl)

    tech_cols = [
        ("Frontend", [
            "Flutter 3.x — Cross-platform UI framework",
            "Dart — Strongly-typed programming language",
            "Riverpod — Reactive state management",
            "GoRouter — Declarative routing & RBAC redirect",
            "Cached Network Image — Efficient image loading",
        ]),
        ("Backend & Data", [
            "Supabase — BaaS platform (PostgreSQL + Auth + Storage + Realtime)",
            "PostgreSQL — Relational database with RLS",
            "Supabase Edge Functions — Deno serverless logic",
            "Supabase Storage — Photo and document storage",
            "Supabase Realtime — WebSocket live sync",
        ]),
        ("Architecture & Tools", [
            "Clean Architecture — domain / data / presentation layers",
            "OneSignal — Push notification service",
            "AES-256 CBC — National ID encryption",
            "flutter_secure_storage — Secure key storage",
            "Android Studio / VS Code — Development IDEs",
        ]),
    ]

    xs = [Inches(0.25), Inches(4.55), Inches(8.85)]
    w  = Inches(4.0)
    for i, (title, pts) in enumerate(tech_cols):
        x = xs[i]
        section_card(sl, title, pts, x, Inches(1.55), w, Inches(5.35), LIGHT_BG)

    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — IMPLEMENTATION HIGHLIGHTS
# ══════════════════════════════════════════════════════════════════════════════
def slide_implementation():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Implementation Highlights", "Key engineering decisions")
    gold_line(sl)

    highlights = [
        "Real-Time Database Integration — Supabase Realtime pushes live bid updates without polling; StreamProvider.family drives UI reactivity",
        "Role-Based Routing — GoRouter redirect resolves user role on every navigation, with 5-second cache to eliminate startup lag",
        "Atomic Bid Placement — Edge Function validates bid amount server-side, updates auction, and notifies outbid users in a single transaction",
        "Dynamic Dashboards — Stat counters update instantly using FutureProvider with refreshable state; no manual refresh required",
        "Image Upload Pipeline — Supabase Storage with UUID-named paths; Cached Network Image prevents redundant downloads",
        "Region-Based Filtering — All auction queries scoped by region_id; admins only see their own auctions and bids",
        "National ID Security — AES-256 CBC encryption applied before any write; key stored in flutter_secure_storage, never in DB",
        "Responsive Android Layouts — All screens validated on Android 7–14; overflow resolved with Flexible + SingleChildScrollView pattern",
    ]
    bullet_list(sl, highlights, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.3), size=14)
    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CHALLENGES FACED
# ══════════════════════════════════════════════════════════════════════════════
def slide_challenges():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Challenges & Solutions", "Real obstacles encountered during development")
    gold_line(sl)

    challenges = [
        ("RenderFlex Overflow",         "Multiple screens overflowed on smaller screens → Replaced Row/Column with Flexible + SingleChildScrollView"),
        ("Supabase RLS Policies",       "Row-level security blocked legitimate queries → Debugged and rewrote RLS rules per role and operation"),
        ("Password Reset UX",           "Reset sheet showed hidden snackbar errors → Replaced with inline error text widget inside the modal"),
        ("Android 7 Compatibility",     "Black screen on startup due to slow async router → Added 5s role cache + 15s global timeout"),
        ("Wireless ADB Testing",        "No USB port access → Used Android Wireless Debugging (adb connect) throughout development"),
        ("State Synchronization",       "Bid counts and auction status fell out of sync → Used StreamProvider.family with Realtime subscriptions"),
        ("Dynamic Count Updates",       "Dashboard stats were static after actions → Switched to FutureProvider with invalidate-on-action pattern"),
    ]

    for i, (challenge, solution) in enumerate(challenges):
        y = Inches(1.62) + Inches(0.72) * i
        add_rect(sl, Inches(0.3), y, Inches(3.3), Inches(0.6), RNP_BLUE)
        add_text(sl, challenge,
                 Inches(0.45), y + Inches(0.09), Inches(3.0), Inches(0.42),
                 size=12, bold=True, color=WHITE)
        add_text(sl, "→  " + solution,
                 Inches(3.75), y + Inches(0.09), Inches(9.0), Inches(0.44),
                 size=12, color=DARK_TEXT)

    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LIVE DEMO FLOW
# ══════════════════════════════════════════════════════════════════════════════
def slide_demo():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Live Demo Flow", "Walkthrough sequence for the demonstration")
    gold_line(sl)

    steps = [
        ("1", "User Login",           "Open app → Enter credentials → Role resolved → Navigate to Home"),
        ("2", "Region Selection",     "Client selects their district/region on first login"),
        ("3", "Browse Auctions",      "Home screen shows active auctions in real time for the selected region"),
        ("4", "Place a Bid",          "Open auction detail → Enter bid amount → Confirm → Bid submitted via Edge Function"),
        ("5", "Notifications",        "Demonstrate outbid push notification arriving on the device"),
        ("6", "Admin Dashboard",      "Login as Region Admin → View live auction stats and active bids"),
        ("7", "Close Auction",        "Admin closes auction → Winner announced → Status updated across all clients"),
        ("8", "Reports",              "Admin generates and downloads region auction PDF report"),
        ("9", "Super Admin View",     "Login as Super Admin → National dashboard → Manage region admins"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(1.55) + Inches(0.585) * i
        add_rect(sl, Inches(0.3), y, Inches(0.5), Inches(0.5), GOLD)
        add_text(sl, num,
                 Inches(0.3), y + Inches(0.05), Inches(0.5), Inches(0.4),
                 size=15, bold=True, color=RNP_BLUE, align=PP_ALIGN.CENTER)
        add_text(sl, title,
                 Inches(0.95), y + Inches(0.07), Inches(2.5), Inches(0.38),
                 size=13, bold=True, color=RNP_BLUE)
        add_text(sl, desc,
                 Inches(3.6), y + Inches(0.07), Inches(9.3), Inches(0.38),
                 size=12, color=DARK_TEXT)

    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ACHIEVEMENTS
# ══════════════════════════════════════════════════════════════════════════════
def slide_achievements():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Achievements", "What we successfully delivered")
    gold_line(sl)

    items_l = [
        "Fully working Flutter Android application",
        "Real-time auction management and bidding system",
        "Secure three-role authentication with GoRouter RBAC",
        "Complete admin and super admin dashboards",
        "Push notifications via OneSignal integration",
        "Dynamic PDF report generation per region",
    ]
    items_r = [
        "AES-256 National ID encryption in secure storage",
        "Supabase Edge Functions for all critical business logic",
        "Production-quality Clean Architecture codebase",
        "Region-based filtering and scoped admin access",
        "Responsive layouts tested on Android 7 through 14",
        "Successful end-to-end testing via wireless ADB",
    ]
    two_col_bullets(sl, items_l, items_r, y=Inches(1.65), size=14)
    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — FUTURE IMPROVEMENTS
# ══════════════════════════════════════════════════════════════════════════════
def slide_future():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Future Improvements", "Planned enhancements beyond this release")
    gold_line(sl)

    future = [
        ("Mobile Money Integration",
         ["MTN Mobile Money & Airtel payment gateway", "Direct bid deposit and refund automation"]),
        ("SMS Notifications",
         ["Twilio/Africa's Talking SMS fallback", "Reach users without push notification enabled"]),
        ("AI-Powered Price Prediction",
         ["ML model trained on historical auction data", "Suggests fair bid ranges for new vehicles"]),
        ("Email Verification",
         ["OTP email confirmation on registration", "Prevents fake account creation"]),
        ("Expanded Auction Categories",
         ["Beyond vehicles: equipment, furniture, assets", "Category-based filtering and browsing"]),
        ("Kinyarwanda Localization",
         ["Full app translation to Kinyarwanda", "Accessible to Rwandan citizens across literacy levels"]),
    ]

    positions = [
        (Inches(0.3),  Inches(1.55), Inches(3.9), Inches(2.5)),
        (Inches(4.5),  Inches(1.55), Inches(3.9), Inches(2.5)),
        (Inches(8.7),  Inches(1.55), Inches(4.2), Inches(2.5)),
        (Inches(0.3),  Inches(4.25), Inches(3.9), Inches(2.5)),
        (Inches(4.5),  Inches(4.25), Inches(3.9), Inches(2.5)),
        (Inches(8.7),  Inches(4.25), Inches(4.2), Inches(2.5)),
    ]
    for i, (title, pts) in enumerate(future):
        x, y, w, h = positions[i]
        section_card(sl, title, pts, x, y, w, h)

    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
def slide_conclusion():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header_bar(sl, "Conclusion", "Summary and impact")
    gold_line(sl)

    add_text(sl,
        "E-CYAMUNARA successfully demonstrates how mobile technology can transform public-sector "
        "auction management in Rwanda.",
        Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.8),
        size=16, color=DARK_TEXT)

    conclusions = [
        "Digitized the full auction lifecycle — from posting to winner announcement — eliminating paper-based processes",
        "Empowered citizens to participate in RNP auctions remotely, increasing transparency and reach",
        "Delivered a production-quality Flutter application with Clean Architecture, Riverpod, and Supabase",
        "Implemented strong security: AES-256 encryption, Supabase RLS, server-side Edge Functions, and RBAC",
        "Demonstrated real value for Rwanda National Police: faster auctions, better records, national reporting",
        "Contributes to Rwanda's digital transformation agenda and the broader e-government ecosystem",
        "Scalable architecture ready for future enhancements including mobile money and AI-assisted pricing",
    ]
    bullet_list(sl, conclusions, Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.4), size=15)
    footer(sl, SLIDE_NUM[0])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
def slide_thankyou():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDE_NUM[0] += 1
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, RNP_BLUE)
    add_rect(sl, 0, Inches(3.5), SLIDE_W, Inches(0.12), GOLD)

    add_text(sl, "Thank You",
             Inches(0.6), Inches(1.0), Inches(12), Inches(1.6),
             size=62, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "Questions & Answers",
             Inches(0.6), Inches(2.55), Inches(12), Inches(0.85),
             size=30, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(sl, "RNP E-CYAMUNARA  ·  Online Vehicle Auction Platform",
             Inches(0.6), Inches(3.85), Inches(12), Inches(0.5),
             size=15, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "GROUP 6  ·  Mobile Application and Handheld Devices  ·  Academic Year 2024–2025",
             Inches(0.6), Inches(4.4), Inches(12), Inches(0.45),
             size=13, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

    members_line = "IRUMVA Jean Marie   ·   MANZI Stiven   ·   NDAYISHIMIYE Issa   ·   NIYOMUGABO Eduard"
    add_text(sl, members_line,
             Inches(0.6), Inches(5.1), Inches(12), Inches(0.45),
             size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # No footer number on final slide – just the brand bar
    add_rect(sl, 0, Inches(7.15), SLIDE_W, Inches(0.35), GOLD)
    add_text(sl, "RP – Musanze Campus  |  2025",
             Inches(0.3), Inches(7.16), Inches(12.7), Inches(0.3),
             size=9, bold=True, color=RNP_BLUE, align=PP_ALIGN.CENTER)

# ── Build all slides ───────────────────────────────────────────────────────────
slide_title()
slide_introduction()
slide_problem()
slide_solution()
slide_roles()
slide_features()
slide_tech()
slide_implementation()
slide_challenges()
slide_demo()
slide_achievements()
slide_future()
slide_conclusion()
slide_thankyou()

# ── Save ───────────────────────────────────────────────────────────────────────
out = r"D:\LEVEL 3\Semester II\MOBILE APPLICATION\ecyamunara\E-CYAMUNARA_PRESENTATION.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {SLIDE_NUM[0]}")
