"""
E-CYAMUNARA Final Presentation Builder
RNP Blue #003087 | Gold #FFD700
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import os, copy
from lxml import etree

# ── Constants ─────────────────────────────────────────────────────────────────
RNP_BLUE  = RGBColor(0x00, 0x30, 0x87)
GOLD      = RGBColor(0xFF, 0xD7, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF9)
MID_GRAY  = RGBColor(0x55, 0x65, 0x7A)
ACCENT    = RGBColor(0x00, 0x6A, 0xD4)

W = Inches(13.33)   # Widescreen 16:9 width
H = Inches(7.5)     # Widescreen 16:9 height

BASE_DIR   = r"D:\LEVEL 3\Semester II\MOBILE APPLICATION\ecyamunara"
RNP_LOGO   = os.path.join(BASE_DIR, "assets", "images", "RNP_logo.png")
PLAT_LOGO  = os.path.join(BASE_DIR, "assets", "images", "PLATFORM-LOGO.png")
OUT_FILE   = os.path.join(BASE_DIR, "E-CYAMUNARA_FINAL_PRESENTATION.pptx")
FOOTER_TXT = "RNP E-CYAMUNARA  |  Online Vehicle Auction Platform"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # Completely blank

# ── Helper: solid fill on shape ────────────────────────────────────────────────
def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_fill(shape):
    shape.fill.background()

def no_line(shape):
    shape.line.fill.background()

# ── Helper: add rectangle ─────────────────────────────────────────────────────
def rect(slide, l, t, w, h, color, line_color=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill_solid(shp, color)
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    else:
        no_line(shp)
    return shp

# ── Helper: add text box ───────────────────────────────────────────────────────
def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box

# ── Helper: multi-paragraph text box ──────────────────────────────────────────
def txbox_lines(slide, lines, l, t, w, h,
                size=16, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, line_spacing_pt=None):
    """lines = list of (text, size, bold, color) or just strings."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, sz, bd, clr = item, size, bold, color
        else:
            txt = item[0]
            sz  = item[1] if len(item)>1 else size
            bd  = item[2] if len(item)>2 else bold
            clr = item[3] if len(item)>3 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(sz)
        run.font.bold  = bd
        run.font.color.rgb = clr
    return box

# ── Helper: footer ─────────────────────────────────────────────────────────────
def add_footer(slide, slide_num=None):
    # Bottom bar
    bar = rect(slide, 0, Inches(7.1), W, Inches(0.4), RNP_BLUE)
    # Footer text
    ftxt = FOOTER_TXT
    if slide_num:
        ftxt += f"   |   Slide {slide_num}"
    txbox(slide, ftxt,
          Inches(0.2), Inches(7.12), Inches(11), Inches(0.35),
          size=10, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

# ── Helper: add logo ───────────────────────────────────────────────────────────
def add_logo(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)

# ── Helper: section header strip ──────────────────────────────────────────────
def section_bar(slide, title):
    rect(slide, 0, 0, W, Inches(1.1), RNP_BLUE)
    rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
    txbox(slide, title,
          Inches(0.35), Inches(0.15), Inches(10), Inches(0.85),
          size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# ── Helper: bullet block ───────────────────────────────────────────────────────
def bullet_block(slide, items, l, t, w, h,
                 size=17, color=DARK_GRAY, bullet="▪"):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if isinstance(item, tuple):
            run.text = f"{bullet}  {item[0]}"
            run.font.size  = Pt(item[1] if len(item)>1 else size)
            run.font.bold  = item[2] if len(item)>2 else False
            run.font.color.rgb = item[3] if len(item)>3 else color
        else:
            run.text = f"{bullet}  {item}"
            run.font.size  = Pt(size)
            run.font.color.rgb = color
    return box

# ── Helper: card / panel ───────────────────────────────────────────────────────
def panel(slide, title, items, l, t, w, h,
          bg=LIGHT_BG, title_color=RNP_BLUE, item_color=DARK_GRAY,
          title_size=16, item_size=14):
    rect(slide, l, t, w, h, bg)
    # Title strip inside card
    rect(slide, l, t, w, Inches(0.38), RNP_BLUE)
    txbox(slide, title,
          l+Inches(0.08), t+Inches(0.04), w-Inches(0.16), Inches(0.32),
          size=title_size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    bullet_block(slide, items,
                 l+Inches(0.1), t+Inches(0.44), w-Inches(0.2), h-Inches(0.5),
                 size=item_size, color=item_color)

# ── Helper: screenshot placeholder ────────────────────────────────────────────
def screenshot_ph(slide, label, l, t, w, h):
    r = rect(slide, l, t, w, h, LIGHT_BG, RNP_BLUE)
    # Dashed effect via line
    r.line.color.rgb = RNP_BLUE
    r.line.width = Pt(1.5)
    txbox(slide, "[ INSERT SCREENSHOT HERE ]",
          l, t+h/2-Inches(0.4), w, Inches(0.5),
          size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    txbox(slide, label,
          l, t+h-Inches(0.4), w, Inches(0.38),
          size=12, bold=False, color=RNP_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
def slide_title():
    sl = prs.slides.add_slide(blank_layout)
    # Background
    rect(sl, 0, 0, W, H, DARK_GRAY)
    # Blue left panel
    rect(sl, 0, 0, Inches(4.8), H, RNP_BLUE)
    # Gold accent strip
    rect(sl, Inches(4.8), 0, Inches(0.07), H, GOLD)

    # Logos on blue panel
    add_logo(sl, RNP_LOGO,  Inches(0.3),  Inches(0.3), Inches(1.4), Inches(1.4))
    add_logo(sl, PLAT_LOGO, Inches(2.9),  Inches(0.3), Inches(1.4), Inches(1.4))

    txbox(sl, "REPUBLIC OF RWANDA",
          Inches(0.2), Inches(1.85), Inches(4.4), Inches(0.4),
          size=11, bold=False, color=GOLD, align=PP_ALIGN.CENTER)
    txbox(sl, "RWANDA NATIONAL POLICE",
          Inches(0.2), Inches(2.2), Inches(4.4), Inches(0.4),
          size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Divider
    rect(sl, Inches(0.4), Inches(2.65), Inches(4.0), Inches(0.04), GOLD)

    txbox(sl, "E-CYAMUNARA",
          Inches(0.2), Inches(2.8), Inches(4.4), Inches(0.85),
          size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, "Online Vehicle Auction Platform",
          Inches(0.2), Inches(3.6), Inches(4.4), Inches(0.45),
          size=16, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

    rect(sl, Inches(0.4), Inches(4.1), Inches(4.0), Inches(0.04), GOLD)

    inst_lines = [
        ("University of Rwanda", 13, True, WHITE),
        ("College of Science and Technology", 11, False, GOLD),
        ("Dept. of Information Technology", 11, False, WHITE),
        ("Academic Year: 2024 – 2025", 11, False, GOLD),
    ]
    txbox_lines(sl, inst_lines,
                Inches(0.2), Inches(4.2), Inches(4.4), Inches(1.2),
                align=PP_ALIGN.CENTER)

    # Right panel — team + course
    txbox(sl, "PROJECT PRESENTATION",
          Inches(5.1), Inches(0.4), Inches(7.8), Inches(0.55),
          size=14, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    txbox(sl, "Mobile Application Development",
          Inches(5.1), Inches(0.9), Inches(7.8), Inches(0.42),
          size=19, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    rect(sl, Inches(5.1), Inches(1.38), Inches(7.8), Inches(0.04), GOLD)

    txbox(sl, "Team Members",
          Inches(5.1), Inches(1.5), Inches(7.8), Inches(0.38),
          size=14, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    members = [
        "NDAYISHIMIYE Issa            —  ID: 222019872",
        "KARANGWA Jean Paul          —  ID: [ID]",
        "UWIMANA Consolee             —  ID: [ID]",
        "NKURUNZIZA Thierry           —  ID: [ID]",
    ]
    bullet_block(sl, members,
                 Inches(5.1), Inches(1.95), Inches(7.8), Inches(1.4),
                 size=14, color=WHITE, bullet="")

    rect(sl, Inches(5.1), Inches(3.45), Inches(7.8), Inches(0.04), GOLD)

    meta = [
        ("Course/Module:", 13, True, GOLD),
        ("Mobile Application Development  (INFT 3204)", 13, False, WHITE),
        ("", 8, False, WHITE),
        ("Supervisor / Lecturer:", 13, True, GOLD),
        ("[Lecturer Name]  —  University of Rwanda", 13, False, WHITE),
        ("", 8, False, WHITE),
        ("Presentation Date:", 13, True, GOLD),
        ("May 2025", 13, False, WHITE),
    ]
    txbox_lines(sl, meta,
                Inches(5.1), Inches(3.55), Inches(7.8), Inches(2.8),
                align=PP_ALIGN.LEFT)

    add_footer(sl, 1)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════════════════════
def slide_toc():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Presentation Outline")

    col1 = [
        "01  Introduction & Background",
        "02  Problem Statement",
        "03  Project Objectives",
        "04  Proposed Solution",
        "05  Requirement Analysis",
        "06  System Architecture",
        "07  Technology Stack",
        "08  UI/UX Design",
        "09  Database Design",
    ]
    col2 = [
        "10  Implementation",
        "11  Testing",
        "12  Challenges & Solutions",
        "13  Live Demo",
        "14  Impact & Benefits",
        "15  Future Improvements",
        "16  Conclusion",
        "17  Questions & Answers",
    ]

    for i, item in enumerate(col1):
        y = Inches(1.3) + i * Inches(0.58)
        rect(sl, Inches(0.4), y, Inches(0.05), Inches(0.4), GOLD)
        txbox(sl, item, Inches(0.6), y, Inches(5.8), Inches(0.45),
              size=15, color=DARK_GRAY)

    for i, item in enumerate(col2):
        y = Inches(1.3) + i * Inches(0.58)
        rect(sl, Inches(6.9), y, Inches(0.05), Inches(0.4), GOLD)
        txbox(sl, item, Inches(7.1), y, Inches(5.8), Inches(0.45),
              size=15, color=DARK_GRAY)

    add_footer(sl, 2)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — INTRODUCTION / BACKGROUND
# ═══════════════════════════════════════════════════════════════════════════════
def slide_intro():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Introduction & Background")

    # Left card
    panel(sl, "Project Context",
          ["E-CYAMUNARA is a Flutter-based mobile auction platform developed for the Rwanda National Police (RNP) to digitize the vehicle auction process.",
           "Vehicle auctions are a core government activity — vehicles seized or retired by RNP require transparent, accessible, and accountable disposal.",
           "The platform replaces cumbersome paper-based processes with a real-time digital solution accessible to all Rwandan citizens."],
          Inches(0.3), Inches(1.25), Inches(6.2), Inches(2.5),
          item_size=13)

    panel(sl, "Digital Rwanda Context",
          ["Rwanda's Vision 2050 prioritises a knowledge-based, digitally-enabled economy.",
           "The government's Smart Rwanda Master Plan targets e-government adoption across all public institutions.",
           "Digital auction platforms reduce corruption, increase public trust, and align with Rwanda's Digital Transformation ambitions.",
           "RNP's adoption of E-CYAMUNARA directly supports the national digitisation mandate."],
          Inches(6.7), Inches(1.25), Inches(6.3), Inches(2.5),
          item_size=13)

    panel(sl, "Why This Project Matters",
          ["Improves transparency in government asset disposal",
           "Enables remote nationwide participation in auctions",
           "Provides real-time bidding and instant notifications",
           "Generates auditable reports for accountability",
           "Reduces operational costs and manual workload for RNP"],
          Inches(0.3), Inches(3.95), Inches(12.7), Inches(2.7),
          item_size=13)

    add_footer(sl, 3)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
def slide_problem():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Problem Statement")

    problems = [
        ("Manual / Paper-Based Processes", "All auction records, bids, and communications managed on paper — slow, error-prone, and costly."),
        ("Lack of Transparency",           "No audit trail, no public visibility into bid prices or results — opens the door to manipulation."),
        ("Limited Accessibility",          "Citizens must physically travel to RNP offices to participate — excludes those in remote regions."),
        ("Poor Bid Tracking",              "No real-time visibility into current highest bid; participants cannot track their own bids."),
        ("Delayed Communication",          "Winners notified days later by mail or phone — no instant confirmation or outbid alerts."),
        ("Fraud Risks",                    "Paper bids can be altered, lost, or fabricated; no cryptographic verification of submissions."),
    ]

    cols = 3
    card_w = Inches(4.2)
    card_h = Inches(2.3)
    gap_x  = Inches(0.15)
    gap_y  = Inches(0.15)
    start_x = Inches(0.3)
    start_y = Inches(1.3)

    for idx, (title, desc) in enumerate(problems):
        col = idx % cols
        row = idx // cols
        lx = start_x + col * (card_w + gap_x)
        ty = start_y + row * (card_h + gap_y)
        rect(sl, lx, ty, card_w, card_h, WHITE)
        rect(sl, lx, ty, Inches(0.06), card_h, RNP_BLUE)
        rect(sl, lx, ty, card_w, Inches(0.04), GOLD)
        txbox(sl, title,
              lx+Inches(0.12), ty+Inches(0.08), card_w-Inches(0.2), Inches(0.4),
              size=14, bold=True, color=RNP_BLUE)
        txbox(sl, desc,
              lx+Inches(0.12), ty+Inches(0.52), card_w-Inches(0.2), Inches(1.65),
              size=12, color=MID_GRAY)

    add_footer(sl, 4)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — PROJECT OBJECTIVES
# ═══════════════════════════════════════════════════════════════════════════════
def slide_objectives():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Project Objectives")

    # General objective
    rect(sl, Inches(0.3), Inches(1.3), Inches(12.7), Inches(0.9), RNP_BLUE)
    txbox(sl, "General Objective",
          Inches(0.5), Inches(1.32), Inches(3.0), Inches(0.4),
          size=13, bold=True, color=GOLD)
    txbox(sl, "To design, develop, and deploy a secure, real-time mobile auction platform for Rwanda National Police that automates vehicle auction management, enhances transparency, and enables nationwide citizen participation.",
          Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.5),
          size=13, color=WHITE)

    txbox(sl, "Specific Objectives",
          Inches(0.3), Inches(2.32), Inches(6.0), Inches(0.38),
          size=16, bold=True, color=RNP_BLUE)
    rect(sl, Inches(0.3), Inches(2.68), Inches(5.8), Inches(0.04), GOLD)

    spec = [
        ("SO1", "Design a secure, role-based mobile application for clients, region admins, and super admins"),
        ("SO2", "Implement real-time bidding using Supabase Realtime with instant outbid notifications"),
        ("SO3", "Implement AES-256 encryption for sensitive personal data (National IDs)"),
        ("SO4", "Develop Edge Functions for atomic bid placement and auction lifecycle management"),
        ("SO5", "Generate PDF auction reports with region-based filtering for administrative oversight"),
        ("SO6", "Integrate push notifications via OneSignal for real-time participant communication"),
        ("SO7", "Provide a feedback mechanism for continuous system improvement"),
        ("SO8", "Ensure Android compatibility across API levels 24–34 (Android 7.0 – 14)"),
    ]

    for i, (code, obj) in enumerate(spec):
        col = i % 2
        row = i // 2
        lx = Inches(0.3) + col * Inches(6.55)
        ty = Inches(2.78) + row * Inches(0.92)
        rect(sl, lx, ty, Inches(6.3), Inches(0.82), WHITE)
        rect(sl, lx, ty, Inches(0.55), Inches(0.82), RNP_BLUE)
        txbox(sl, code,
              lx+Inches(0.04), ty+Inches(0.2), Inches(0.48), Inches(0.4),
              size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(sl, obj,
              lx+Inches(0.62), ty+Inches(0.08), Inches(5.6), Inches(0.68),
              size=12, color=DARK_GRAY)

    add_footer(sl, 5)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — PROPOSED SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
def slide_solution():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Proposed Solution — E-CYAMUNARA")

    # Overview
    txbox(sl, "E-CYAMUNARA is a cross-platform mobile application that digitises the complete auction lifecycle — from vehicle posting and bid submission to auction closure and winner notification — all in real time.",
          Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.7),
          size=14, color=MID_GRAY, italic=True)

    flow_steps = [
        ("1", "Admin Posts\nAuction"),
        ("2", "Citizens\nRegister"),
        ("3", "View &\nBid Online"),
        ("4", "Real-Time\nUpdates"),
        ("5", "Auction\nClosed"),
        ("6", "Winner\nNotified"),
    ]

    step_w  = Inches(1.9)
    step_h  = Inches(1.0)
    gap     = Inches(0.22)
    start_x = Inches(0.3)
    ty      = Inches(2.05)

    for i, (num, label) in enumerate(flow_steps):
        lx = start_x + i * (step_w + gap)
        rect(sl, lx, ty, step_w, step_h, RNP_BLUE)
        txbox(sl, num,
              lx, ty+Inches(0.05), step_w, Inches(0.35),
              size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txbox(sl, label,
              lx, ty+Inches(0.38), step_w, Inches(0.58),
              size=11, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(flow_steps)-1:
            txbox(sl, "→",
                  lx+step_w+Inches(0.04), ty+Inches(0.25), Inches(0.18), Inches(0.5),
                  size=18, bold=True, color=RNP_BLUE)

    # Three pillars
    pillars = [
        ("Role-Based Access Control",
         ["Three-tier: Super Admin → Region Admin → Client",
          "Separate dashboards per role",
          "GoRouter redirect guards on every navigation"]),
        ("Real-Time Bidding Engine",
         ["Supabase Realtime subscription on auction rows",
          "Atomic bid placement via Edge Functions",
          "Outbid notifications within 1–2 seconds"]),
        ("Transparency & Reporting",
         ["All bids permanently recorded in PostgreSQL",
          "PDF report generation per region / period",
          "Full audit trail — tamper-evident"]),
    ]

    pil_w = Inches(4.1)
    for i, (title, items) in enumerate(pillars):
        lx = Inches(0.3) + i*(pil_w+Inches(0.2))
        panel(sl, title, items,
              lx, Inches(3.25), pil_w, Inches(3.35),
              item_size=12)

    add_footer(sl, 6)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — FUNCTIONAL REQUIREMENTS
# ═══════════════════════════════════════════════════════════════════════════════
def slide_func_req():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Requirement Analysis — Functional Requirements")

    reqs = [
        ("User Management",    ["User registration with NID encryption", "Multi-role login (Client / Admin / Super Admin)", "Account suspension & activation", "Profile management"]),
        ("Auction Management", ["Post auctions with photos & details", "Set start/end dates, starting price", "Region-based auction filtering", "Manual / automatic auction closure"]),
        ("Bidding System",     ["Place bids (min increment enforced)", "Real-time current price updates", "Bid history per auction", "My Bids tracking dashboard"]),
        ("Notifications",      ["Outbid push notification via OneSignal", "Auction start/end alerts", "In-app notification centre", "Admin broadcast messages"]),
        ("Reporting",          ["PDF report generation per region", "National consolidated reports", "Report download & storage", "Admin dashboard statistics"]),
        ("Feedback",           ["Client feedback submission", "Admin view & respond to feedback", "Feedback rating & categorisation"]),
    ]

    card_w = Inches(4.1)
    card_h = Inches(2.45)
    for idx, (title, items) in enumerate(reqs):
        col = idx % 3
        row = idx // 3
        lx = Inches(0.3) + col*(card_w+Inches(0.17))
        ty = Inches(1.3) + row*(card_h+Inches(0.12))
        panel(sl, title, items, lx, ty, card_w, card_h, item_size=12)

    add_footer(sl, 7)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — NON-FUNCTIONAL REQUIREMENTS
# ═══════════════════════════════════════════════════════════════════════════════
def slide_nonfunc_req():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Requirement Analysis — Non-Functional Requirements")

    rows = [
        ("Security",       "AES-256 NID encryption; Supabase RLS policies; JWT authentication; secure storage"),
        ("Performance",    "Bid placement < 2 s; Realtime update < 1.5 s; App startup < 3 s on mid-range devices"),
        ("Scalability",    "Supabase managed PostgreSQL scales horizontally; Edge Functions scale automatically"),
        ("Reliability",    "99.9 % uptime target via Supabase cloud infrastructure; edge function retries"),
        ("Availability",   "24/7 access; offline graceful degradation with clear error messaging"),
        ("Usability",      "Intuitive Material Design UI; WCAG contrast compliance; clear action hierarchy"),
        ("Maintainability","Clean Architecture (data/domain/presentation); single provider import; documented constants"),
        ("Compatibility",  "Android API 24–34 (Android 7.0–14); Flutter stable channel; responsive for 5–7″ screens"),
    ]

    # Table header
    rect(sl, Inches(0.3), Inches(1.28), Inches(3.2), Inches(0.45), RNP_BLUE)
    rect(sl, Inches(3.5), Inches(1.28), Inches(9.5), Inches(0.45), RNP_BLUE)
    txbox(sl, "Quality Attribute", Inches(0.35), Inches(1.3), Inches(3.1), Inches(0.4),
          size=13, bold=True, color=WHITE)
    txbox(sl, "Requirement / Acceptance Criterion", Inches(3.55), Inches(1.3), Inches(9.4), Inches(0.4),
          size=13, bold=True, color=WHITE)

    for i, (attr, desc) in enumerate(rows):
        ty = Inches(1.73) + i * Inches(0.6)
        bg = WHITE if i % 2 == 0 else LIGHT_BG
        rect(sl, Inches(0.3), ty, Inches(3.2), Inches(0.58), bg)
        rect(sl, Inches(3.5), ty, Inches(9.5), Inches(0.58), bg)
        txbox(sl, attr, Inches(0.4), ty+Inches(0.08), Inches(3.0), Inches(0.42),
              size=13, bold=True, color=RNP_BLUE)
        txbox(sl, desc, Inches(3.6), ty+Inches(0.08), Inches(9.3), Inches(0.42),
              size=12, color=DARK_GRAY)

    add_footer(sl, 8)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — SYSTEM ARCHITECTURE (diagram-style)
# ═══════════════════════════════════════════════════════════════════════════════
def slide_architecture():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, DARK_GRAY)
    # Header
    rect(sl, 0, 0, W, Inches(1.0), RNP_BLUE)
    rect(sl, 0, Inches(1.0), W, Inches(0.06), GOLD)
    txbox(sl, "System Architecture",
          Inches(0.35), Inches(0.12), Inches(10), Inches(0.78),
          size=30, bold=True, color=WHITE)

    # LAYER 1 — Flutter Client
    rect(sl, Inches(0.2), Inches(1.2), Inches(3.6), Inches(2.6), RGBColor(0,72,148))
    txbox(sl, "FLUTTER CLIENT", Inches(0.25), Inches(1.25), Inches(3.5), Inches(0.38),
          size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    for i, lbl in enumerate(["Presentation Layer\n(Screens / Widgets)",
                               "Riverpod State\nManagement",
                               "GoRouter Navigation\n& RBAC Guards"]):
        ty = Inches(1.7) + i * Inches(0.7)
        rect(sl, Inches(0.35), ty, Inches(3.3), Inches(0.58), RGBColor(0,48,120))
        txbox(sl, lbl, Inches(0.4), ty+Inches(0.04), Inches(3.2), Inches(0.52),
              size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow →
    txbox(sl, "HTTPS\nREST / WS", Inches(3.85), Inches(2.2), Inches(1.1), Inches(0.55),
          size=10, color=GOLD, align=PP_ALIGN.CENTER)
    txbox(sl, "⟶", Inches(3.9), Inches(2.0), Inches(1.0), Inches(0.4),
          size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # LAYER 2 — Supabase
    rect(sl, Inches(5.0), Inches(1.2), Inches(4.2), Inches(5.4), RGBColor(0,72,148))
    txbox(sl, "SUPABASE BACKEND", Inches(5.05), Inches(1.25), Inches(4.1), Inches(0.38),
          size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    supa_items = [
        "Auth (GoTrue / JWT)",
        "PostgreSQL Database",
        "Row Level Security (RLS)",
        "Supabase Realtime",
        "Object Storage",
        "Edge Functions (Deno)",
    ]
    for i, lbl in enumerate(supa_items):
        ty = Inches(1.7) + i * Inches(0.72)
        rect(sl, Inches(5.15), ty, Inches(3.9), Inches(0.6), RGBColor(0,48,120))
        txbox(sl, lbl, Inches(5.2), ty+Inches(0.08), Inches(3.8), Inches(0.44),
              size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow →
    txbox(sl, "⟶", Inches(9.3), Inches(2.3), Inches(0.6), Inches(0.4),
          size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # LAYER 3 — External Services
    rect(sl, Inches(9.95), Inches(1.2), Inches(3.1), Inches(5.4), RGBColor(0,72,148))
    txbox(sl, "EXTERNAL SERVICES", Inches(9.95), Inches(1.25), Inches(3.05), Inches(0.38),
          size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    ext_items = [
        "OneSignal\nPush Notifications",
        "Secure Storage\n(flutter_secure_storage)",
        "PDF Generation\n(Edge Function)",
        "Android OS\nAPI 24–34",
    ]
    for i, lbl in enumerate(ext_items):
        ty = Inches(1.7) + i * Inches(1.2)
        rect(sl, Inches(10.05), ty, Inches(2.9), Inches(1.0), RGBColor(0,48,120))
        txbox(sl, lbl, Inches(10.1), ty+Inches(0.08), Inches(2.8), Inches(0.84),
              size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Three-role indicators at bottom
    roles = ["CLIENT", "REGION ADMIN", "SUPER ADMIN"]
    role_colors = [RGBColor(0,106,212), RGBColor(0,160,120), RGBColor(200,80,0)]
    for i, (role, col) in enumerate(zip(roles, role_colors)):
        lx = Inches(0.25) + i * Inches(1.25)
        rect(sl, lx, Inches(3.98), Inches(1.15), Inches(0.42), col)
        txbox(sl, role, lx, Inches(4.0), Inches(1.15), Inches(0.38),
              size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(sl, 9)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════════
def slide_tech():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Technology Stack")

    categories = [
        ("Frontend", [
            "Flutter (Dart) — Cross-platform UI framework",
            "flutter_riverpod — State management",
            "go_router — Declarative navigation & RBAC",
            "flutter_animate — Smooth UI animations",
            "cached_network_image — Optimised image loading",
        ]),
        ("Backend & Database", [
            "Supabase — BaaS platform (PostgreSQL 15)",
            "Row Level Security (RLS) — Fine-grained access control",
            "Supabase Realtime — WebSocket-based live updates",
            "Deno Edge Functions — Serverless business logic",
            "Supabase Storage — Auction & profile photos",
        ]),
        ("Security & Services", [
            "flutter_secure_storage — AES-256 key vault",
            "AES-CBC Encryption — National ID protection",
            "Supabase Auth — JWT + GoTrue authentication",
            "OneSignal — Cross-platform push notifications",
            "pdf — Server-side PDF report generation",
        ]),
        ("Dev Tools", [
            "VS Code + Flutter extension",
            "Android Studio + Emulator",
            "Wireless ADB (adb connect) — Cable-free testing",
            "GitHub — Version control & collaboration",
            "Supabase CLI — Local dev & function deploy",
        ]),
    ]

    card_w = Inches(6.3)
    card_h = Inches(2.75)
    for idx, (cat, items) in enumerate(categories):
        col = idx % 2
        row = idx // 2
        lx = Inches(0.3) + col*(card_w+Inches(0.4))
        ty = Inches(1.3) + row*(card_h+Inches(0.12))
        panel(sl, cat, items, lx, ty, card_w, card_h, item_size=13)

    add_footer(sl, 10)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — UI/UX DESIGN PHILOSOPHY
# ═══════════════════════════════════════════════════════════════════════════════
def slide_uiux_design():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "UI/UX Design — Design Philosophy & Theme")

    # Color swatches
    swatches = [
        (RNP_BLUE,              "#003087",  "RNP Blue\n(Primary)"),
        (GOLD,                  "#FFD700",  "Gold\n(Accent)"),
        (RGBColor(26,26,46),    "#1A1A2E",  "Dark\n(Background)"),
        (RGBColor(244,246,249), "#F4F6F9",  "Light Grey\n(Surfaces)"),
        (WHITE,                             "#FFFFFF",  "White\n(Cards)"),
    ]
    tx = Inches(0.3)
    for i, (color, hex_c, label) in enumerate(swatches):
        lx = tx + i*Inches(2.55)
        rect(sl, lx, Inches(1.3), Inches(2.2), Inches(0.75), color)
        txbox(sl, hex_c, lx, Inches(2.08), Inches(2.2), Inches(0.3),
              size=11, bold=True, color=RNP_BLUE, align=PP_ALIGN.CENTER)
        txbox(sl, label, lx, Inches(2.38), Inches(2.2), Inches(0.5),
              size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Design principles
    principles = [
        ("Material Design 3",  "Follows Google Material You guidelines for familiarity and consistency across Android devices."),
        ("Role-Aware UI",      "Each of the three roles (Client / Region Admin / Super Admin) sees a tailored interface and navigation bar."),
        ("Contrast Compliance","RNP Blue + White + Gold palette meets WCAG AA contrast standards for readability."),
        ("Responsive Layout",  "Adaptive to 5–7 inch screens; SingleChildScrollView and flexible widgets prevent overflow."),
        ("Projection-Ready",   "High-contrast colours, large touch targets (min 44dp), and clear typography work on projectors."),
        ("Android 7 Support",  "Tested on API 24+; no features restricted to higher APIs without fallback."),
    ]

    for i, (title, desc) in enumerate(principles):
        col = i % 2
        row = i // 2
        lx = Inches(0.3) + col*Inches(6.55)
        ty = Inches(3.0) + row*Inches(1.22)
        rect(sl, lx, ty, Inches(6.3), Inches(1.12), WHITE)
        rect(sl, lx, ty, Inches(0.06), Inches(1.12), GOLD)
        txbox(sl, title, lx+Inches(0.12), ty+Inches(0.06), Inches(6.0), Inches(0.38),
              size=13, bold=True, color=RNP_BLUE)
        txbox(sl, desc, lx+Inches(0.12), ty+Inches(0.48), Inches(6.0), Inches(0.6),
              size=12, color=MID_GRAY)

    add_footer(sl, 11)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — UI/UX SCREENSHOTS (placeholders)
# ═══════════════════════════════════════════════════════════════════════════════
def slide_screenshots():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "UI/UX Design — Application Screens")

    screens = [
        ("Splash Screen",       Inches(0.3),  Inches(1.25)),
        ("Login Screen",        Inches(2.8),  Inches(1.25)),
        ("Registration",        Inches(5.3),  Inches(1.25)),
        ("Client Home",         Inches(7.8),  Inches(1.25)),
        ("Auction Detail",      Inches(10.3), Inches(1.25)),
        ("My Bids",             Inches(0.3),  Inches(4.15)),
        ("Notifications",       Inches(2.8),  Inches(4.15)),
        ("Region Admin Dash",   Inches(5.3),  Inches(4.15)),
        ("Super Admin Dash",    Inches(7.8),  Inches(4.15)),
        ("Reports",             Inches(10.3), Inches(4.15)),
    ]

    ph_w = Inches(2.3)
    ph_h = Inches(2.7)
    for label, lx, ty in screens:
        screenshot_ph(sl, label, lx, ty, ph_w, ph_h)

    add_footer(sl, 12)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — DATABASE DESIGN
# ═══════════════════════════════════════════════════════════════════════════════
def slide_database():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Database Design — Entity Overview")

    tables = [
        ("users",         ["id (uuid PK)", "full_name", "email", "national_id (encrypted)", "phone", "role", "is_suspended", "profile_photo_url", "created_at"]),
        ("auctions",      ["id (uuid PK)", "title", "description", "category", "region", "starting_price", "current_price", "status", "start_date", "end_date", "admin_id (FK)", "photos[]"]),
        ("bids",          ["id (uuid PK)", "auction_id (FK)", "user_id (FK)", "amount", "created_at"]),
        ("notifications", ["id (uuid PK)", "user_id (FK)", "title", "body", "type", "is_read", "created_at"]),
        ("feedback",      ["id (uuid PK)", "user_id (FK)", "message", "rating", "status", "created_at"]),
        ("region_admins", ["id (uuid PK)", "full_name", "email", "region", "is_active", "created_at"]),
    ]

    col1 = tables[:3]
    col2 = tables[3:]

    card_w = Inches(6.3)
    gap_y  = Inches(0.12)
    card_h = Inches(1.9)
    for col_idx, col_tables in enumerate([col1, col2]):
        lx = Inches(0.3) + col_idx*(card_w+Inches(0.4))
        for row_idx, (tname, fields) in enumerate(col_tables):
            ty = Inches(1.3) + row_idx*(card_h+gap_y)
            rect(sl, lx, ty, card_w, card_h, WHITE)
            rect(sl, lx, ty, card_w, Inches(0.38), RNP_BLUE)
            txbox(sl, f"TABLE: {tname.upper()}",
                  lx+Inches(0.08), ty+Inches(0.04), card_w-Inches(0.16), Inches(0.32),
                  size=13, bold=True, color=WHITE)
            # fields in two mini-cols
            half = len(fields)//2 + len(fields)%2
            for fi, field in enumerate(fields):
                fx = lx+Inches(0.1) + (fi//half)*Inches(3.1)
                fy = ty+Inches(0.44) + (fi%half)*Inches(0.27)
                txbox(sl, f"• {field}", fx, fy, Inches(3.0), Inches(0.26),
                      size=10, color=MID_GRAY)

    # Relationships note
    rect(sl, Inches(0.3), Inches(6.92), Inches(12.7), Inches(0.55), RNP_BLUE)
    txbox(sl, "Key Relationships: bids.auction_id → auctions.id  |  bids.user_id → users.id  |  notifications.user_id → users.id  |  auctions.admin_id → region_admins.id",
          Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.48),
          size=11, color=WHITE)

    add_footer(sl, 13)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — IMPLEMENTATION OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
def slide_impl_overview():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Implementation — Key Modules")

    modules = [
        ("Authentication Module",
         ["Multi-role login via Supabase Auth (GoTrue)",
          "AES-256 CBC encryption on National ID before storage",
          "Account suspension check on every session resume",
          "Secure token storage via flutter_secure_storage"]),
        ("Auction Management",
         ["Region admins post auctions with multi-photo upload",
          "Supabase Storage bucket: auction-photos",
          "Auction lifecycle: OPEN → CLOSING → CLOSED",
          "Super admin can manage all auctions nationally"]),
        ("Real-Time Bidding",
         ["StreamProvider.family backed by Supabase Realtime",
          "place-bid Edge Function: atomic amount validation",
          "Current price updated instantly across all clients",
          "BidNotifier (StateNotifier) manages placement state"]),
        ("Notifications System",
         ["OneSignal SDK for cross-device push delivery",
          "send-auction-notification Edge Function dispatches alerts",
          "In-app notification list with read/unread state",
          "Outbid, auction start/end, and admin broadcast types"]),
        ("Reporting Engine",
         ["generate-report Edge Function creates PDF with dart:pdf",
          "Region-filtered and national-level report modes",
          "Reports saved to Supabase Storage: reports bucket",
          "Download link returned to admin after generation"]),
        ("Feedback System",
         ["Clients submit feedback via FeedbackScreen",
          "Admins view and categorise feedback submissions",
          "Rating system 1–5 stars with open-ended message",
          "Stored in feedback table with status lifecycle"]),
    ]

    card_w = Inches(4.1)
    card_h = Inches(2.7)
    for idx, (title, items) in enumerate(modules):
        col = idx % 3
        row = idx // 3
        lx = Inches(0.3) + col*(card_w+Inches(0.17))
        ty = Inches(1.3) + row*(card_h+Inches(0.1))
        panel(sl, title, items, lx, ty, card_w, card_h, item_size=12)

    add_footer(sl, 14)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — IMPLEMENTATION SCREENSHOTS
# ═══════════════════════════════════════════════════════════════════════════════
def slide_impl_screens():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Implementation — Screen Gallery")

    screens = [
        ("Login / Auth Flow",       Inches(0.3),  Inches(1.25), Inches(3.1), Inches(2.8)),
        ("Client Home Dashboard",   Inches(3.55), Inches(1.25), Inches(3.1), Inches(2.8)),
        ("Auction Detail & Bid",    Inches(6.8),  Inches(1.25), Inches(3.1), Inches(2.8)),
        ("Admin Auction Management",Inches(10.05),Inches(1.25), Inches(3.1), Inches(2.8)),
        ("Bid Confirmation Screen", Inches(0.3),  Inches(4.25), Inches(3.1), Inches(2.65)),
        ("Notifications Centre",    Inches(3.55), Inches(4.25), Inches(3.1), Inches(2.65)),
        ("Reports Dashboard",       Inches(6.8),  Inches(4.25), Inches(3.1), Inches(2.65)),
        ("Super Admin Panel",       Inches(10.05),Inches(4.25), Inches(3.1), Inches(2.65)),
    ]
    for label, lx, ty, pw, ph in screens:
        screenshot_ph(sl, label, lx, ty, pw, ph)

    add_footer(sl, 15)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — TESTING
# ═══════════════════════════════════════════════════════════════════════════════
def slide_testing():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Testing & Quality Assurance")

    test_areas = [
        ("Functional Testing",
         ["All user stories verified against acceptance criteria",
          "CRUD operations for auctions, bids, users, feedback",
          "Role-based access control validated per route"]),
        ("Real-Time Testing",
         ["Supabase Realtime events verified across 3 concurrent sessions",
          "Bid synchronisation latency measured < 1.5 s",
          "No missed events under standard network conditions"]),
        ("Authentication Testing",
         ["Login / logout / session persistence tested",
          "Suspended account block verified on navigation",
          "Token expiry and refresh flow validated"]),
        ("Android Compatibility",
         ["Tested on API 24 (Android 7), 28 (9), 30 (11), 33 (13), 34 (14)",
          "Wireless ADB used for physical device testing",
          "No API-version-specific crashes detected"]),
        ("UI / Responsiveness",
         ["Screen sizes from 5″ to 7″ tested",
          "No widget overflow on smallest supported screen",
          "Scroll behaviour verified on all list views"]),
        ("Edge Function Testing",
         ["place-bid: concurrent bid race conditions tested",
          "close-auction: winner assignment verified atomically",
          "create-admin: RLS bypass tested with service role"]),
    ]

    card_w = Inches(4.1)
    card_h = Inches(2.55)
    for idx, (title, items) in enumerate(test_areas):
        col = idx % 3
        row = idx // 3
        lx = Inches(0.3) + col*(card_w+Inches(0.17))
        ty = Inches(1.3) + row*(card_h+Inches(0.1))
        panel(sl, title, items, lx, ty, card_w, card_h, item_size=12)

    add_footer(sl, 16)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — CHALLENGES & SOLUTIONS
# ═══════════════════════════════════════════════════════════════════════════════
def slide_challenges():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Challenges Encountered & Solutions Applied")

    rows = [
        ("RLS Policy Conflicts",     "Supabase Row Level Security blocked legitimate admin reads on certain tables.",          "Designed explicit policies per role; used Edge Functions with service role key for admin operations."),
        ("Black Screen on Startup",  "GoRouter redirect made 3 sequential Supabase queries per navigation, causing a 5–10 s black screen on Android 7.",  "Implemented 5-second role cache + 15-second query timeout inside the redirect callback."),
        ("Android 7 Compatibility",  "Certain Flutter plugins and Material 3 behaviours unsupported below API 28.",           "Pinned minSdkVersion to 24; replaced incompatible widgets; tested on physical API-24 device via Wireless ADB."),
        ("Concurrent Bid Race Cond.","Two users placing identical bids simultaneously caused inconsistent highest-bid state.", "Moved bid logic to place-bid Edge Function using PostgreSQL transaction for atomic compare-and-update."),
        ("Realtime Sync Lag",        "Supabase Realtime events occasionally delayed under poor connectivity.",                "Added optimistic UI update on bid submission; Realtime confirms or corrects within 1–2 s."),
        ("Widget Overflow",          "Long auction titles and descriptions caused RenderFlex overflow on small screens.",      "Wrapped all list items in Flexible/Expanded; added TextOverflow.ellipsis with tooltip on tap."),
        ("Notification Delivery",    "OneSignal tokens not registered before first login, causing missed push notifications.", "Added notification initialisation inside post-login flow; fall back to in-app notification centre."),
        ("State After Suspension",   "Suspended users could remain in session until token expiry.",                           "Added suspension check in GoRouter redirect; every navigation triggers a lightweight status query."),
    ]

    # Header row
    rect(sl, Inches(0.3), Inches(1.3), Inches(3.8), Inches(0.42), RNP_BLUE)
    rect(sl, Inches(4.15), Inches(1.3), Inches(4.2), Inches(0.42), RNP_BLUE)
    rect(sl, Inches(8.4), Inches(1.3), Inches(4.6), Inches(0.42), RNP_BLUE)
    for lx, w, lbl in [(Inches(0.35), Inches(3.7), "Challenge"),
                        (Inches(4.2),  Inches(4.1), "Root Cause"),
                        (Inches(8.45), Inches(4.5), "Solution")]:
        txbox(sl, lbl, lx, Inches(1.33), w, Inches(0.36),
              size=13, bold=True, color=WHITE)

    for i, (ch, cause, sol) in enumerate(rows):
        ty = Inches(1.75) + i * Inches(0.67)
        bg = WHITE if i%2==0 else LIGHT_BG
        rect(sl, Inches(0.3), ty, Inches(3.8), Inches(0.65), bg)
        rect(sl, Inches(4.15), ty, Inches(4.2), Inches(0.65), bg)
        rect(sl, Inches(8.4), ty, Inches(4.6), Inches(0.65), bg)
        txbox(sl, ch,    Inches(0.35), ty+Inches(0.05), Inches(3.7), Inches(0.58),
              size=11, bold=True, color=RNP_BLUE)
        txbox(sl, cause, Inches(4.2), ty+Inches(0.05), Inches(4.1), Inches(0.58),
              size=11, color=MID_GRAY)
        txbox(sl, sol,   Inches(8.45), ty+Inches(0.05), Inches(4.5), Inches(0.58),
              size=11, color=DARK_GRAY)

    add_footer(sl, 17)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — LIVE DEMO FLOW
# ═══════════════════════════════════════════════════════════════════════════════
def slide_demo():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, DARK_GRAY)
    rect(sl, 0, 0, W, Inches(1.0), RNP_BLUE)
    rect(sl, 0, Inches(1.0), W, Inches(0.06), GOLD)
    txbox(sl, "Live Demo — Walkthrough Flow",
          Inches(0.35), Inches(0.12), Inches(10), Inches(0.78),
          size=30, bold=True, color=WHITE)

    demo_steps = [
        ("01", "Splash &\nApp Launch"),
        ("02", "Client\nRegister"),
        ("03", "Login &\nRole Detect"),
        ("04", "Region\nSelection"),
        ("05", "Client\nDashboard"),
        ("06", "View\nAuction"),
        ("07", "Place\nBid"),
        ("08", "Real-Time\nUpdate"),
        ("09", "Notification\nReceived"),
        ("10", "Admin\nDashboard"),
        ("11", "Close\nAuction"),
        ("12", "View\nReports"),
    ]

    step_w = Inches(0.98)
    step_h = Inches(1.35)
    gap    = Inches(0.07)
    start_x= Inches(0.3)
    row1_y = Inches(1.25)
    row2_y = Inches(2.85)

    for i, (num, label) in enumerate(demo_steps):
        col = i % 6
        row_y = row1_y if i < 6 else row2_y
        lx = start_x + col*(step_w+gap)
        rect(sl, lx, row_y, step_w, step_h, RGBColor(0,48,120))
        rect(sl, lx, row_y, step_w, Inches(0.38), RNP_BLUE)
        txbox(sl, num, lx, row_y+Inches(0.03), step_w, Inches(0.34),
              size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txbox(sl, label, lx, row_y+Inches(0.42), step_w, Inches(0.88),
              size=11, color=WHITE, align=PP_ALIGN.CENTER)
        if col < 5:
            txbox(sl, "→", lx+step_w+Inches(0.0), row_y+Inches(0.55),
                  Inches(0.08), Inches(0.35),
                  size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Presenter notes area
    rect(sl, Inches(0.3), Inches(4.4), Inches(12.7), Inches(2.35), RGBColor(10,30,70))
    txbox(sl, "PRESENTER NOTES",
          Inches(0.45), Inches(4.45), Inches(4.0), Inches(0.38),
          size=12, bold=True, color=GOLD)
    notes = [
        "• Start with the app already installed on a physical device or emulator.",
        "• Use two devices or browser tabs to demonstrate real-time bid synchronisation.",
        "• Show the admin closing an auction and the client receiving the winner notification.",
        "• Open a report PDF to demonstrate the reporting engine. Keep demo under 15 minutes.",
    ]
    for i, note in enumerate(notes):
        txbox(sl, note,
              Inches(0.45), Inches(4.88)+i*Inches(0.44), Inches(12.4), Inches(0.42),
              size=12, color=WHITE)

    add_footer(sl, 18)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — IMPACT & BENEFITS
# ═══════════════════════════════════════════════════════════════════════════════
def slide_impact():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Impact & Benefits")

    beneficiaries = [
        ("Rwanda National Police",
         ["Eliminates manual paperwork and physical record storage",
          "Real-time dashboard visibility into all national auctions",
          "Automated PDF reports for accountability and audit",
          "Reduced staff workload through system automation",
          "Traceable and tamper-evident bid records"]),
        ("Rwandan Citizens",
         ["Nationwide remote participation — no travel required",
          "Fair and transparent bidding with equal access",
          "Instant push notifications on bid status changes",
          "Easy-to-use mobile interface for all literacy levels",
          "Secure handling of personal identification data"]),
        ("Digital Rwanda Agenda",
         ["Demonstrates feasibility of mobile e-government services",
          "Aligns with Vision 2050 digital transformation goals",
          "Open architecture reusable for other govt. auctions",
          "Reduces carbon footprint by eliminating physical presence",
          "Contributes to Rwanda's global ICT competitiveness"]),
    ]

    card_w = Inches(4.1)
    for i, (title, items) in enumerate(beneficiaries):
        lx = Inches(0.3) + i*(card_w+Inches(0.2))
        panel(sl, title, items,
              lx, Inches(1.3), card_w, Inches(5.3),
              item_size=13)

    add_footer(sl, 19)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — FUTURE IMPROVEMENTS
# ═══════════════════════════════════════════════════════════════════════════════
def slide_future():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, LIGHT_BG)
    section_bar(sl, "Future Improvements & Roadmap")

    short_term = [
        "iOS / iPadOS support via Flutter platform adaptation",
        "Kinyarwanda language localisation (i18n)",
        "SMS fallback notifications via Twilio / Rwanda MTN",
        "Biometric authentication (fingerprint / face ID)",
        "Offline-first mode with local Hive cache & sync",
    ]
    mid_term = [
        "Mobile Money integration (MoMo API, Airtel Money)",
        "Web admin dashboard (Flutter Web or Next.js)",
        "Advanced analytics with charts & trend analysis",
        "AI-powered price prediction per vehicle category",
        "Bulk auction import from Excel / CSV",
    ]
    long_term = [
        "Machine-learning bid prediction for reserve pricing",
        "AI fraud detection flagging anomalous bid patterns",
        "National cloud scaling across all government agencies",
        "Open API for third-party auction integrations",
        "Cross-border auction support for EAC member states",
    ]

    panel(sl, "Short-Term (0–6 months)", short_term,
          Inches(0.3), Inches(1.3), Inches(4.1), Inches(4.75), item_size=13)
    panel(sl, "Mid-Term (6–18 months)", mid_term,
          Inches(4.6), Inches(1.3), Inches(4.1), Inches(4.75), item_size=13)
    panel(sl, "Long-Term (18+ months)", long_term,
          Inches(8.9), Inches(1.3), Inches(4.1), Inches(4.75), item_size=13)

    add_footer(sl, 20)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
def slide_conclusion():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, DARK_GRAY)
    rect(sl, 0, 0, W, Inches(1.0), RNP_BLUE)
    rect(sl, 0, Inches(1.0), W, Inches(0.06), GOLD)
    txbox(sl, "Conclusion",
          Inches(0.35), Inches(0.12), Inches(10), Inches(0.78),
          size=30, bold=True, color=WHITE)

    # Summary statement
    rect(sl, Inches(0.3), Inches(1.2), Inches(12.7), Inches(1.0), RGBColor(0,48,120))
    txbox(sl, "E-CYAMUNARA successfully delivers a production-ready, secure, and scalable mobile auction platform that modernises the Rwanda National Police vehicle disposal process and sets a replicable model for digital government services in Rwanda.",
          Inches(0.5), Inches(1.28), Inches(12.3), Inches(0.84),
          size=14, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    achievements = [
        ("Technical Achievements",
         ["Clean Architecture with Riverpod state management",
          "AES-256 NID encryption + Supabase RLS security model",
          "Atomic real-time bidding via Deno Edge Functions",
          "PDF report generation and cloud storage integration",
          "Android 7.0–14 compatibility verified on physical devices"]),
        ("Academic Contributions",
         ["Demonstrates full-stack Flutter + Supabase integration",
          "Applies RBAC design in a real government use case",
          "Implements serverless functions for tamper-proof logic",
          "Contributes a replicable e-government mobile app pattern",
          "Complete documentation of architecture and design decisions"]),
        ("Social / Business Value",
         ["Increases transparency in government asset disposal",
          "Enables nationwide equitable access to RNP auctions",
          "Reduces operational costs and human error for RNP",
          "Aligns with Rwanda Vision 2050 digital economy goals",
          "Provides a foundation for future AI and ML enhancements"]),
    ]

    card_w = Inches(4.1)
    for i, (title, items) in enumerate(achievements):
        lx = Inches(0.3) + i*(card_w+Inches(0.2))
        panel(sl, title, items,
              lx, Inches(2.4), card_w, Inches(4.2),
              bg=RGBColor(10,30,70), title_color=GOLD,
              item_color=WHITE, item_size=12)

    add_footer(sl, 21)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 22 — THANK YOU / Q&A
# ═══════════════════════════════════════════════════════════════════════════════
def slide_qa():
    sl = prs.slides.add_slide(blank_layout)
    rect(sl, 0, 0, W, H, DARK_GRAY)
    rect(sl, 0, 0, Inches(5.0), H, RNP_BLUE)
    rect(sl, Inches(5.0), 0, Inches(0.07), H, GOLD)

    add_logo(sl, RNP_LOGO,  Inches(0.5), Inches(0.4), Inches(1.6), Inches(1.6))
    add_logo(sl, PLAT_LOGO, Inches(2.9), Inches(0.4), Inches(1.6), Inches(1.6))

    txbox(sl, "RWANDA NATIONAL POLICE",
          Inches(0.1), Inches(2.1), Inches(4.8), Inches(0.42),
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, "E-CYAMUNARA",
          Inches(0.1), Inches(2.6), Inches(4.8), Inches(0.85),
          size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, "Online Vehicle Auction Platform",
          Inches(0.1), Inches(3.42), Inches(4.8), Inches(0.45),
          size=15, color=GOLD, align=PP_ALIGN.CENTER)
    rect(sl, Inches(0.5), Inches(3.95), Inches(4.0), Inches(0.05), GOLD)
    txbox(sl, "University of Rwanda  |  Academic Year 2024–2025",
          Inches(0.1), Inches(4.08), Inches(4.8), Inches(0.38),
          size=12, color=WHITE, align=PP_ALIGN.CENTER)

    txbox(sl, "Thank You",
          Inches(5.3), Inches(1.3), Inches(7.7), Inches(0.9),
          size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    rect(sl, Inches(5.3), Inches(2.22), Inches(7.4), Inches(0.05), GOLD)
    txbox(sl, "Questions & Answers",
          Inches(5.3), Inches(2.38), Inches(7.7), Inches(0.6),
          size=26, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

    txbox(sl, "Team Members",
          Inches(5.3), Inches(3.2), Inches(7.7), Inches(0.38),
          size=14, bold=True, color=GOLD)

    members = [
        "NDAYISHIMIYE Issa",
        "KARANGWA Jean Paul",
        "UWIMANA Consolee",
        "NKURUNZIZA Thierry",
    ]
    for i, m in enumerate(members):
        txbox(sl, f"▪  {m}",
              Inches(5.3), Inches(3.65)+i*Inches(0.5), Inches(7.4), Inches(0.44),
              size=14, color=WHITE)

    txbox(sl, "Mobile Application Development  —  INFT 3204",
          Inches(5.3), Inches(5.9), Inches(7.7), Inches(0.4),
          size=13, color=MID_GRAY, italic=True)

    add_footer(sl, 22)

# ═══════════════════════════════════════════════════════════════════════════════
#  BUILD ALL SLIDES
# ═══════════════════════════════════════════════════════════════════════════════
slide_title()
slide_toc()
slide_intro()
slide_problem()
slide_objectives()
slide_solution()
slide_func_req()
slide_nonfunc_req()
slide_architecture()
slide_tech()
slide_uiux_design()
slide_screenshots()
slide_database()
slide_impl_overview()
slide_impl_screens()
slide_testing()
slide_challenges()
slide_demo()
slide_impact()
slide_future()
slide_conclusion()
slide_qa()

prs.save(OUT_FILE)
print(f"Saved: {OUT_FILE}")
print(f"Total slides: {len(prs.slides)}")
